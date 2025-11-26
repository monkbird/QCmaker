from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os
from typing import List, Dict, Any

class PPTService:
    def generate_ppt(self, project_name: str, topic: str, data_summary: str, discussion_summary: str, chart_images: List[str] = []) -> str:
        prs = Presentation()
        
        # Slide 1: Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = project_name
        subtitle.text = f"QC Topic: {topic}"
        
        # Slide 2: Project Overview
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "项目概况 (Project Overview)"
        tf = body_shape.text_frame
        tf.text = "背景与数据汇总"
        
        p = tf.add_paragraph()
        p.text = data_summary
        p.level = 1

        # Slide 3: Discussion & Analysis
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "分析与讨论 (Analysis & Discussion)"
        tf = body_shape.text_frame
        tf.text = discussion_summary

        # Slide 4: Conclusion
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "结论 (Conclusion)"
        tf = body_shape.text_frame
        tf.text = "根据上述分析，我们建议..."

        # Save
        output_dir = "backend/generated_ppts"
        os.makedirs(output_dir, exist_ok=True)
        filename = f"{project_name}_QC_Report.pptx"
        output_path = os.path.join(output_dir, filename)
        prs.save(output_path)
        
        return output_path

ppt_service = PPTService()
